#!/usr/bin/env python3
"""
WR KB Selective Indexer — Phase 3 of the WaterRoads migration

Walks the WaterRoads KB Shared Drive, extracts text from high-signal documents,
embeds via Voyage AI (voyage-3.5, 1024 dim), and inserts into the WR Supabase
documents table. Stores drive_file_id + drive_modified for incremental sync.

Selectivity rules:
    INDEXABLE extensions: .pdf .docx .xlsx .pptx .md .txt .csv
                         + Google native (.gdoc .gsheet .gslides via export)
    SKIP if size > MAX_FILE_BYTES (default 200 MB)
    SKIP if path contains any SKIP_PATH_FRAGMENTS (Pictures, Logos, Recordings, ...)
    SKIP MIME types: image/*, video/*, audio/*, application/zip, etc.

Idempotency:
    - For each Drive file, check WR documents.drive_file_id
    - If exists AND drive_modified matches, skip
    - If exists AND drive_modified differs, delete old + reinsert (re-embed)
    - If not exists, insert

Usage:
    source scripts/env-setup.sh && source .secrets/wr-env.sh

    # Dry run — list what would be indexed without writes
    python3 scripts/wr-index-drive-content.py --dry-run

    # Index just the SharePoint-imported content (smaller, finished first)
    python3 scripts/wr-index-drive-content.py --root "Imported from SharePoint"

    # Index a specific subfolder for testing
    python3 scripts/wr-index-drive-content.py --root "Imported from SharePoint/LGG Water Roads 2025/Diligence Docs"

    # Limit to first N files (useful for first run)
    python3 scripts/wr-index-drive-content.py --limit 20

    # Full run (no limit)
    python3 scripts/wr-index-drive-content.py
"""
import argparse
import io
import json
import os
import re
import sys
import time
from collections import Counter
from datetime import datetime
from pathlib import Path

import httpx
import voyageai
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload


# ---------- Config ----------
MAX_FILE_BYTES = 200 * 1024 * 1024     # 200 MB hard cap
MAX_CHARS_PER_CHUNK = 8000
EMBED_RATE_LIMIT_S = 0.2               # Voyage rate limiting

INDEXABLE_MIMES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-excel.sheet.macroenabled.12": "xlsx",
    "text/markdown": "text",
    "text/plain": "text",
    "text/csv": "csv",
    "application/vnd.google-apps.document": "gdoc",
    "application/vnd.google-apps.spreadsheet": "gsheet",
    "application/vnd.google-apps.presentation": "gslides",
}

SKIP_PATH_FRAGMENTS = [
    "/Pictures/", "/Photos/", "/Logos/", "/Recordings/", "/Camera Uploads/",
    "/Screenshots/", "/Renders/", "/Stock Imagery/",
]


# ---------- Drive helpers ----------
def get_drive_service():
    sa = os.environ["WR_SERVICE_ACCOUNT_FILE"]
    creds = service_account.Credentials.from_service_account_file(
        sa, scopes=["https://www.googleapis.com/auth/drive"])
    return build("drive", "v3", credentials=creds)


def find_folder_by_path(svc, drive_id, path):
    """path like 'Imported from SharePoint/LGG Water Roads 2025'. Returns folder id."""
    parent = drive_id
    for part in [p for p in path.split("/") if p]:
        q = (f"name = '{part.replace(chr(39), chr(39)+chr(39))}' "
             f"and '{parent}' in parents and trashed = false "
             f"and mimeType = 'application/vnd.google-apps.folder'")
        res = svc.files().list(q=q, supportsAllDrives=True, includeItemsFromAllDrives=True,
            driveId=drive_id, corpora="drive", fields="files(id,name)").execute()
        files = res.get("files", [])
        if not files:
            return None
        parent = files[0]["id"]
    return parent


def walk_drive(svc, drive_id, folder_id, prefix=""):
    """Yield (full_path, file_record) for every non-folder file under folder_id."""
    page = None
    while True:
        res = svc.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
            driveId=drive_id, corpora="drive",
            fields="nextPageToken,files(id,name,mimeType,size,modifiedTime,webViewLink)",
            pageSize=200, pageToken=page,
        ).execute()
        for f in res.get("files", []):
            full = f"{prefix}/{f['name']}" if prefix else f['name']
            if f["mimeType"] == "application/vnd.google-apps.folder":
                yield from walk_drive(svc, drive_id, f["id"], full)
            else:
                yield full, f
        page = res.get("nextPageToken")
        if not page: break


# ---------- Text extraction ----------
def extract_pdf(data: bytes) -> str:
    try:
        import pdfplumber
    except ImportError:
        return ""
    out = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages[:300]:  # cap at 300 pages
            try:
                out.append(page.extract_text() or "")
            except Exception:
                pass
    return "\n\n".join(out)


def extract_docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""
    doc = Document(io.BytesIO(data))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also tables
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paras.append(" | ".join(cells))
    return "\n".join(paras)


def extract_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    out = []
    for sheet in wb.worksheets:
        out.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(max_rows=2000, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                out.append(" | ".join(cells))
    return "\n".join(out)


def extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides):
        out.append(f"## Slide {i+1}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip()
                    if txt:
                        out.append(txt)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"[notes] {notes}")
    return "\n".join(out)


def extract_text(data: bytes) -> str:
    try:
        return data.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def extract_via_export(svc, file_id, export_mime):
    """For Google native files — export then extract."""
    req = svc.files().export_media(fileId=file_id, mimeType=export_mime)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    data = buf.getvalue()
    if export_mime == "text/plain":
        return extract_text(data)
    if export_mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_docx(data)
    return ""


def download_and_extract(svc, file_record):
    fid = file_record["id"]
    mime = file_record["mimeType"]
    kind = INDEXABLE_MIMES.get(mime)

    if kind in ("gdoc", "gsheet", "gslides"):
        # Export Google native files
        if kind == "gdoc":
            return extract_via_export(svc, fid, "text/plain")
        if kind == "gsheet":
            return extract_via_export(svc, fid, "text/csv")
        if kind == "gslides":
            return extract_via_export(svc, fid, "text/plain")

    # Binary download for office/text formats
    req = svc.files().get_media(fileId=fid, supportsAllDrives=True)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    data = buf.getvalue()

    if kind == "pdf":     return extract_pdf(data)
    if kind == "docx":    return extract_docx(data)
    if kind == "xlsx":    return extract_xlsx(data)
    if kind == "pptx":    return extract_pptx(data)
    if kind in ("text", "csv"): return extract_text(data)
    return ""


# ---------- Categorisation ----------
def categorise(path: str) -> str:
    """Derive category from the top-level meaningful folder.

    For 'Imported from SharePoint/Diligence Docs/foo.pdf' → 'diligence_docs'.
    For 'Governance/Board Papers/Approved/foo.docx'      → 'governance'.
    """
    parts = [p for p in path.split("/") if p]
    if not parts: return "general"
    # Strip the "Imported from X" wrapper but keep the next folder
    if parts[0].startswith("Imported from") and len(parts) > 1:
        parts = parts[1:]
    if len(parts) < 2:
        # File at top-level — no folder context
        return "general"
    base = parts[0].lower()
    base = re.sub(r"^\d+[\.\-_\s]*", "", base)
    base = re.sub(r"[^a-z0-9]+", "_", base).strip("_")
    return base or "general"


# ---------- Supabase helpers ----------
def supa_url(): return os.environ["WR_SUPABASE_URL"]
def supa_headers():
    k = os.environ["WR_SUPABASE_SERVICE_ROLE_KEY"]
    return {"apikey": k, "Authorization": f"Bearer {k}", "Content-Type": "application/json"}


def already_indexed(drive_file_id, drive_modified):
    r = httpx.get(f"{supa_url()}/rest/v1/documents", headers=supa_headers(),
        params={"drive_file_id": f"eq.{drive_file_id}",
                "select": "id,drive_modified", "limit": 1}, timeout=15)
    if r.status_code == 200 and r.json():
        existing = r.json()[0]
        if existing.get("drive_modified") == drive_modified:
            return "skip"  # unchanged
        return "update"  # re-embed
    return "new"


def delete_existing(drive_file_id):
    httpx.delete(f"{supa_url()}/rest/v1/documents?drive_file_id=eq.{drive_file_id}",
        headers=supa_headers(), timeout=15)


def chunk_text(text, max_chars=MAX_CHARS_PER_CHUNK):
    text = text.strip()
    if len(text) <= max_chars:
        return [text] if text else []
    chunks = []
    for i in range(0, len(text), max_chars):
        chunks.append(text[i:i+max_chars])
    return chunks


def insert_chunks(vo, file_record, full_path, text, category):
    if not text or not text.strip():
        return 0
    chunks = chunk_text(text)
    title = Path(file_record["name"]).stem
    inserted = 0
    for i, chunk in enumerate(chunks):
        try:
            embedding = vo.embed([chunk], model="voyage-3.5", input_type="document").embeddings[0]
        except Exception as e:
            print(f"      embed failed chunk {i}: {e}")
            continue
        record = {
            "entity": "waterroads",
            "source_file": full_path,
            "title": f"{title} (Part {i+1})" if len(chunks) > 1 else title,
            "content": chunk,
            "embedding": embedding,
            "category": category,
            "drive_file_id": file_record["id"],
            "drive_modified": file_record.get("modifiedTime"),
            "metadata": {
                "chunk_index": i,
                "total_chunks": len(chunks),
                "drive_url": file_record.get("webViewLink"),
                "mime_type": file_record["mimeType"],
                "size_bytes": int(file_record.get("size", 0) or 0),
                "embedding_model": "voyage-3.5",
            },
        }
        r = httpx.post(f"{supa_url()}/rest/v1/documents",
            headers={**supa_headers(), "Prefer": "return=minimal"},
            json=record, timeout=30)
        if r.status_code in (200, 201):
            inserted += 1
        else:
            print(f"      insert failed chunk {i}: {r.status_code} {r.text[:200]}")
        time.sleep(EMBED_RATE_LIMIT_S)
    return inserted


# ---------- Main ----------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--root", default="",
                    help="Subpath under WaterRoads KB to index (default: whole drive)")
    ap.add_argument("--limit", type=int, default=0,
                    help="Stop after N files (0 = no limit)")
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    print(f"=== WR KB Indexer — {datetime.now().isoformat()} ===")
    print(f"Root: {args.root or '(whole drive)'}, limit: {args.limit or 'none'}, "
          f"dry-run: {args.dry_run}\n")

    svc = get_drive_service()
    drive_id = os.environ["WR_DRIVE_ID"]

    if args.root:
        start_id = find_folder_by_path(svc, drive_id, args.root)
        if not start_id:
            sys.exit(f"FATAL: folder not found: {args.root}")
    else:
        start_id = drive_id

    vo = None
    if not args.dry_run:
        vo = voyageai.Client(api_key=os.environ["VOYAGE_API_KEY"])

    counters = Counter()
    new_inserted = 0
    files_processed = 0

    for full_path, f in walk_drive(svc, drive_id, start_id, args.root):
        if args.limit and files_processed >= args.limit:
            print(f"\n--- limit reached ({args.limit}) ---")
            break

        # Filter checks
        if any(frag in f"/{full_path}/" for frag in SKIP_PATH_FRAGMENTS):
            counters["skip:path_fragment"] += 1
            continue
        size = int(f.get("size", 0) or 0)
        if size > MAX_FILE_BYTES:
            counters["skip:too_big"] += 1
            print(f"  SKIP big ({size//1024//1024}MB): {full_path}")
            continue
        kind = INDEXABLE_MIMES.get(f["mimeType"])
        if not kind:
            counters[f"skip:mime:{f['mimeType']}"] += 1
            continue

        # Idempotency check
        decision = already_indexed(f["id"], f.get("modifiedTime"))
        if decision == "skip":
            counters["skip:unchanged"] += 1
            continue

        files_processed += 1
        category = categorise(full_path)
        size_mb = size / 1024 / 1024

        if args.dry_run:
            print(f"  WOULD INDEX [{decision}] [{category}] ({size_mb:.1f}MB): {full_path}")
            counters[f"would:{decision}"] += 1
            continue

        try:
            print(f"  -> [{category}] ({size_mb:.1f}MB): {full_path}")
            if decision == "update":
                delete_existing(f["id"])
            text = download_and_extract(svc, f)
            n = insert_chunks(vo, f, full_path, text, category)
            new_inserted += n
            counters[f"inserted:{decision}"] += 1
            print(f"      {n} chunks embedded")
        except Exception as e:
            print(f"      ERROR: {e}")
            counters["error"] += 1

    print(f"\n=== Done ===")
    print(f"Files processed: {files_processed}, chunks embedded: {new_inserted}")
    print(f"\nCounters:")
    for k, v in counters.most_common():
        print(f"  {k}: {v}")


if __name__ == "__main__":
    main()
