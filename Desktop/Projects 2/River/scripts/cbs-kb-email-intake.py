#!/usr/bin/env python3
"""
CBS KB Email Intake — add content to the CBS Group knowledge base by email.

How it works:
  1. Jeff or Sarah send an email to jeff@cbs.com.au with subject starting "CBS KB:"
  2. This script scans the inbox via Microsoft Graph for matching emails
  3. Downloads attachments (.pdf, .docx, .xlsx, .pptx, .md, .txt)
  4. Uploads to WR Google Drive (Shared Drive) in the target folder
  5. Extracts text, embeds via Voyage AI, inserts into WR Supabase
  6. Marks the email as read and sends a Teams notification confirming ingestion

Subject line format:
  CBS KB: <category> | <optional title override>

Examples:
  "CBS KB: Governance | Board resolution template 2026"
    → uploads to Governance/ folder, titled "Board resolution template 2026"

  "CBS KB: Financial"
    → uploads to Financial/ folder, uses filename as title

  "CBS KB:"
    → uploads to Incoming/ folder (catch-all)

  "CBS KB: Reference/Shipley | Win Themes Guide"
    → uploads to Reference/Shipley/ folder

Body text (if no attachments):
  If the email has no attachments but has body text, the body itself is indexed
  as a document (useful for quick notes, meeting summaries, observations).

Environment variables:
  MICROSOFT_CLIENT_ID, MICROSOFT_CLIENT_SECRET, MICROSOFT_TENANT_ID
  SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY
  WR_SERVICE_ACCOUNT_FILE, WR_DRIVE_ID
  VOYAGE_API_KEY
  TEAMS_WEBHOOK_URL (optional, for confirmation notifications)

Usage:
  source scripts/env-setup.sh && source .secrets/wr-env.sh
  python3 scripts/wr-kb-email-intake.py              # process new emails
  python3 scripts/wr-kb-email-intake.py --days 7     # look back further
  python3 scripts/wr-kb-email-intake.py --dry-run    # list matches without processing
"""
import argparse
import base64
import io
import json
import os
import re
import sys
import time
from datetime import datetime, timedelta, timezone
from pathlib import Path

import httpx
import msal
import voyageai
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload


# ---------- Config ----------
INBOX_USER = "jeff@cbs.com.au"
SUBJECT_PREFIX = "CBS KB"  # matches "CBS KB:", "CBS KB -", "CBS KB Items" etc.
GRAPH = "https://graph.microsoft.com/v1.0"
MAX_CHARS_PER_CHUNK = 8000
EMBED_RATE_LIMIT_S = 0.25

INDEXABLE_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx", ".md", ".txt", ".csv",
    ".doc", ".xls",  # older formats — best effort
}

FOLDER_MIME = "application/vnd.google-apps.folder"

# URL domains to follow when found in email body
# LinkedIn articles/posts + major news/industry sources
URL_FOLLOW_PATTERNS = [
    r"linkedin\.com/(posts|pulse|feed/update|in/[\w\-]+/recent-activity)",
    r"(reuters|bloomberg|afr|smh|theaustralian|theguardian|abc\.net|bbc)"
    r"\.(com|com\.au|co\.uk|net\.au)",
    r"(infrastructure|transport|water|engineering|government)"
    r"[\w\.\-]*\.(com|com\.au|org|org\.au|gov\.au|govt\.nz)",
]
URL_FOLLOW_RE = [re.compile(p, re.IGNORECASE) for p in URL_FOLLOW_PATTERNS]

# URLs to never follow (images, videos, calendar links, email signatures)
URL_SKIP_PATTERNS = re.compile(
    r"\.(png|jpg|jpeg|gif|svg|mp4|mov|mp3|ics|vcf)(\?|$)|"
    r"(unsubscribe|manage.preferences|email-protection|mailto:|tel:)",
    re.IGNORECASE,
)


# ---------- URL → PDF pipeline ----------
URL_RE_GENERIC = re.compile(r"https?://[^\s\"'<>)\]\,]+", re.IGNORECASE)
FETCH_UA = "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"


def extract_urls(text):
    """Find URLs in email body that should be followed for KB ingestion."""
    if not text: return []
    raw_urls = URL_RE_GENERIC.findall(text)
    follow = []
    seen = set()
    for url in raw_urls:
        url = url.rstrip(".")  # strip trailing punctuation
        if url in seen: continue
        seen.add(url)
        if URL_SKIP_PATTERNS.search(url): continue
        # Check if it matches any follow pattern
        if any(p.search(url) for p in URL_FOLLOW_RE):
            follow.append(url)
        # Also follow ANY url the user explicitly put in a "CBS KB:" email —
        # if they emailed it, they want it indexed
        elif "linkedin.com" in url.lower():
            follow.append(url)
    return follow


def fetch_and_convert_to_pdf(url):
    """Fetch a web page, extract article text, generate a clean PDF.
    Returns (pdf_bytes, extracted_text, page_title) or (None, None, None) on failure."""
    try:
        from bs4 import BeautifulSoup
        from fpdf import FPDF
    except ImportError:
        print(f"    bs4/fpdf2 not installed — cannot fetch URLs")
        return None, None, None

    try:
        r = httpx.get(url, headers={"User-Agent": FETCH_UA}, timeout=30,
                      follow_redirects=True)
        if r.status_code != 200:
            print(f"    fetch {url}: HTTP {r.status_code}")
            return None, None, None
    except Exception as e:
        print(f"    fetch {url}: {e}")
        return None, None, None

    soup = BeautifulSoup(r.text, "html.parser")

    # Extract title
    page_title = ""
    if soup.title and soup.title.string:
        page_title = soup.title.string.strip()
    if not page_title:
        h1 = soup.find("h1")
        page_title = h1.get_text(strip=True) if h1 else url.split("/")[-1][:80]

    # Extract article text — try common article containers first
    article_text = ""
    for selector in ["article", '[role="main"]', ".article-body", ".post-content",
                     ".entry-content", ".story-body", "main"]:
        el = soup.select_one(selector)
        if el:
            paragraphs = el.find_all(["p", "h1", "h2", "h3", "h4", "li", "blockquote"])
            article_text = "\n\n".join(p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True))
            if len(article_text) > 200:
                break

    # Fallback: all paragraphs on page
    if len(article_text) < 200:
        paragraphs = soup.find_all("p")
        article_text = "\n\n".join(p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True))

    if not article_text or len(article_text) < 50:
        print(f"    no meaningful text extracted from {url}")
        return None, None, None

    # Strip to reasonable length (some pages have footer noise)
    article_text = article_text[:50000]

    # Generate PDF — wrapped in try/except so failures don't block the pipeline.
    # If PDF generation fails, we still return the extracted text for embedding.
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_left_margin(15)
        pdf.set_right_margin(15)
        pdf.add_page()

        # Title
        safe_title = page_title[:200].encode("latin-1", errors="replace").decode("latin-1")
        pdf.set_font("Helvetica", "B", 13)
        pdf.multi_cell(0, 7, safe_title)
        pdf.ln(2)

        # Source + capture date — use cell() to avoid wrapping issues with long URLs
        pdf.set_font("Helvetica", "", 7)
        pdf.set_text_color(100, 100, 100)
        source_line = f"Source: {url[:120]}"
        pdf.cell(0, 4, source_line.encode("latin-1", errors="replace").decode("latin-1"), new_x="LMARGIN", new_y="NEXT")
        pdf.cell(0, 4, f"Captured: {datetime.now().strftime('%Y-%m-%d %H:%M')}", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)

        # Article body
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("Helvetica", "", 10)
        for para in article_text.split("\n\n"):
            para = para.strip()
            if not para: continue
            safe = para.encode("latin-1", errors="replace").decode("latin-1")
            pdf.multi_cell(0, 5, safe)
            pdf.ln(2)

        pdf_bytes = bytes(pdf.output())
    except Exception as e:
        print(f"    PDF generation failed ({e}) — will embed text without PDF")
        pdf_bytes = None

    return pdf_bytes, article_text, page_title


# ---------- Auth ----------
def graph_token():
    app = msal.ConfidentialClientApplication(
        os.environ["MICROSOFT_CLIENT_ID"],
        authority=f"https://login.microsoftonline.com/{os.environ['MICROSOFT_TENANT_ID']}",
        client_credential=os.environ["MICROSOFT_CLIENT_SECRET"],
    )
    res = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
    if "access_token" not in res:
        raise RuntimeError(f"Token acquire failed: {res}")
    return res["access_token"]


def get_drive_service():
    sa = os.environ["WR_SERVICE_ACCOUNT_FILE"]
    creds = service_account.Credentials.from_service_account_file(
        sa, scopes=["https://www.googleapis.com/auth/drive"])
    return build("drive", "v3", credentials=creds)


def supa_headers():
    k = os.environ["SUPABASE_SERVICE_ROLE_KEY"]
    return {"apikey": k, "Authorization": f"Bearer {k}", "Content-Type": "application/json"}


def supa_url():
    return os.environ["SUPABASE_URL"]


# ---------- Subject parsing ----------
def parse_subject(subject):
    """Parse 'CBS KB: category | title' or 'CBS KB category' → (category, title_override).
    Tolerant of missing colon — if someone writes 'CBS KB Items' it still works."""
    if not subject: return "incoming", None
    # Strip the prefix (with or without colon)
    rest = subject
    for prefix in ["CBS KB:", "CBS KB -", "CBS KB"]:
        if rest.lower().startswith(prefix.lower()):
            rest = rest[len(prefix):].strip()
            break
    if not rest:
        return "incoming", None
    if "|" in rest:
        cat, title = rest.split("|", 1)
        return cat.strip() or "incoming", title.strip() or None
    return rest.strip(), None


# ---------- Drive folder management ----------
def ensure_folder(svc, drive_id, folder_path):
    """Create/find nested folder path like 'Governance/Board Papers'. Returns folder id."""
    parent = drive_id
    for part in [p.strip() for p in folder_path.split("/") if p.strip()]:
        q = (f"name = '{part}' and '{parent}' in parents "
             f"and mimeType = '{FOLDER_MIME}' and trashed = false")
        res = svc.files().list(q=q, supportsAllDrives=True,
            includeItemsFromAllDrives=True, driveId=drive_id,
            corpora="drive", fields="files(id,name)").execute()
        files = res.get("files", [])
        if files:
            parent = files[0]["id"]
        else:
            body = {"name": part, "mimeType": FOLDER_MIME, "parents": [parent]}
            created = svc.files().create(body=body, supportsAllDrives=True,
                fields="id").execute()
            parent = created["id"]
    return parent


def upload_file(svc, drive_id, folder_id, filename, content_bytes, mime_type):
    """Upload a file to Drive. Returns file metadata dict."""
    media = MediaIoBaseUpload(io.BytesIO(content_bytes), mimetype=mime_type, resumable=False)
    body = {"name": filename, "parents": [folder_id]}
    return svc.files().create(body=body, media_body=media, supportsAllDrives=True,
        fields="id,name,webViewLink,modifiedTime,mimeType,size").execute()


# ---------- Text extraction (reuse from indexer) ----------
def extract_text(data, ext):
    ext = ext.lower()
    if ext == ".pdf":
        try:
            import pdfplumber
            out = []
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                for p in pdf.pages[:200]:
                    try: out.append(p.extract_text() or "")
                    except: pass
            return "\n\n".join(out)
        except ImportError:
            return ""
    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return ""
    elif ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
            out = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(max_rows=2000, values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells: out.append(" | ".join(cells))
            return "\n".join(out)
        except ImportError:
            return ""
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = "".join(r.text for r in para.runs).strip()
                            if txt: out.append(txt)
            return "\n".join(out)
        except ImportError:
            return ""
    elif ext in (".md", ".txt", ".csv"):
        return data.decode("utf-8", errors="ignore")
    return ""


# ---------- Embedding + insert ----------
def chunk_text(text, max_chars=MAX_CHARS_PER_CHUNK):
    text = text.strip()
    if not text: return []
    if len(text) <= max_chars: return [text]
    return [text[i:i+max_chars] for i in range(0, len(text), max_chars)]


def embed_and_insert(vo, text, title, source_file, category, drive_file_id=None, drive_modified=None, email_message_id=None):
    """Embed text chunks and insert into WR Supabase. Returns chunk count."""
    chunks = chunk_text(text)
    if not chunks: return 0
    H = supa_headers()
    inserted = 0
    for i, chunk in enumerate(chunks):
        try:
            embedding = vo.embed([chunk], model="voyage-3.5", input_type="document").embeddings[0]
        except Exception as e:
            print(f"    embed error chunk {i}: {e}")
            continue
        record = {
            "entity": "cbs-group",
            "source_file": source_file,
            "title": f"{title} (Part {i+1})" if len(chunks) > 1 else title,
            "content": chunk,
            "embedding": embedding,
            "category": category.lower().replace(" ", "_").replace("/", "_"),
            "drive_file_id": drive_file_id,
            "drive_modified": drive_modified,
            "metadata": {"chunk_index": i, "total_chunks": len(chunks),
                         "ingestion_method": "email_intake",
                         "email_message_id": email_message_id,
                         "embedding_model": "voyage-3.5"},
        }
        r = httpx.post(f"{supa_url()}/rest/v1/documents", headers=H, json=record, timeout=30)
        if r.status_code in (200, 201):
            inserted += 1
        else:
            print(f"    insert error chunk {i}: {r.status_code} {r.text[:200]}")
        time.sleep(EMBED_RATE_LIMIT_S)
    return inserted


# ---------- Graph email operations ----------
def get_processed_message_ids():
    """Return set of Graph message IDs already processed (stored in documents metadata)."""
    H = supa_headers()
    r = httpx.get(f"{supa_url()}/rest/v1/documents",
        headers=H,
        params={"metadata->>ingestion_method": "eq.email_intake",
                "select": "metadata", "limit": 5000},
        timeout=30)
    ids = set()
    if r.status_code == 200:
        for doc in r.json():
            md = doc.get("metadata") or {}
            mid = md.get("email_message_id")
            if mid: ids.add(mid)
    return ids


def find_kb_emails(token, since_iso):
    """Find emails with subject starting 'CBS KB'. Uses processed ID tracking
    instead of isRead flag (we don't have Mail.ReadWrite to mark emails as read)."""
    H = {"Authorization": f"Bearer {token}", "Prefer": 'outlook.body-content-type="text"'}
    filter_q = (
        f"receivedDateTime ge {since_iso} "
        f"and startsWith(subject, '{SUBJECT_PREFIX}')"
    )
    r = httpx.get(f"{GRAPH}/users/{INBOX_USER}/messages",
        headers=H,
        params={"$filter": filter_q,
                "$select": "id,subject,from,receivedDateTime,body,hasAttachments",
                "$top": 50, "$orderby": "receivedDateTime desc"},
        timeout=30)
    if r.status_code != 200:
        print(f"  Graph search error: {r.status_code} {r.text[:200]}")
        return []
    return r.json().get("value", [])


def get_attachments(token, message_id):
    """Download all file attachments from a message."""
    H = {"Authorization": f"Bearer {token}"}
    r = httpx.get(f"{GRAPH}/users/{INBOX_USER}/messages/{message_id}/attachments",
        headers=H, timeout=60)
    if r.status_code != 200: return []
    out = []
    for att in r.json().get("value", []):
        if att.get("@odata.type") == "#microsoft.graph.fileAttachment":
            out.append({
                "name": att["name"],
                "contentType": att.get("contentType", "application/octet-stream"),
                "data": base64.b64decode(att.get("contentBytes", "")),
            })
    return out


def mark_read(token, message_id):
    """Mark email as read so it's not processed again."""
    H = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    httpx.patch(f"{GRAPH}/users/{INBOX_USER}/messages/{message_id}",
        headers=H, json={"isRead": True}, timeout=15)


# ---------- Teams notification ----------
def notify(title, summary):
    webhook = os.environ.get("TEAMS_WEBHOOK_URL")
    if not webhook: return
    try:
        httpx.post(webhook, json={"title": f"{title}\n{summary}"}, timeout=15)
    except Exception as e:
        print(f"  Teams notify failed: {e}")


# ---------- Main ----------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--days", type=int, default=3)
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    print(f"=== CBS KB Email Intake — {datetime.now(timezone.utc).isoformat()} ===")
    print(f"Lookback: {args.days} days, dry-run: {args.dry_run}")

    token = graph_token()
    since = (datetime.now(timezone.utc) - timedelta(days=args.days)).strftime("%Y-%m-%dT%H:%M:%SZ")
    emails = find_kb_emails(token, since)
    print(f"Found {len(emails)} unread '{SUBJECT_PREFIX}' emails")

    if not emails:
        print("Nothing to process.")
        return

    svc = get_drive_service()
    drive_id = os.environ["WR_DRIVE_ID"]
    vo = None if args.dry_run else voyageai.Client(api_key=os.environ["VOYAGE_API_KEY"])

    total_files = 0
    total_chunks = 0

    # Idempotency: track which emails we've already processed (by Graph message ID)
    processed_ids = set() if args.dry_run else get_processed_message_ids()

    for email in emails:
        msg_id = email.get("id")
        if msg_id in processed_ids:
            print(f"\n  [skip — already processed] {email.get('subject','')[:60]}")
            continue

        subj = email.get("subject", "")
        sender = email.get("from", {}).get("emailAddress", {}).get("address", "")
        category, title_override = parse_subject(subj)
        print(f"\n  [{sender}] {subj}")
        print(f"    → category: {category}, title override: {title_override or '(use filename)'}")

        if args.dry_run:
            if email.get("hasAttachments"):
                atts = get_attachments(token, email["id"])
                for a in atts:
                    print(f"    attachment: {a['name']} ({len(a['data'])/1024:.0f} KB)")
            else:
                body_len = len((email.get("body", {}).get("content", "") or ""))
                print(f"    no attachments, body length: {body_len} chars")
            continue

        folder_id = ensure_folder(svc, drive_id, category)
        body_text = (email.get("body", {}).get("content", "") or "").strip()

        # --- Always index the email body if it has substance ---
        # The body often contains context, instructions, meeting notes, or
        # standalone knowledge — valuable whether or not attachments exist.
        if body_text and len(body_text) > 50:
            body_title = title_override or f"Email: {subj}"
            body_bytes = body_text.encode("utf-8")
            fname = re.sub(r"[^a-zA-Z0-9\-_ ]", "", body_title)[:80] + ".md"
            uploaded_body = upload_file(svc, drive_id, folder_id,
                fname, body_bytes, "text/markdown")
            print(f"    uploaded body as: {fname}")

            n = embed_and_insert(vo, body_text, body_title,
                f"email-intake/{category}/{fname}", category,
                uploaded_body["id"], uploaded_body.get("modifiedTime"),
                email_message_id=msg_id)
            print(f"    body indexed: {n} chunks")
            total_files += 1
            total_chunks += n

        # --- Index each attachment ---
        if email.get("hasAttachments"):
            atts = get_attachments(token, email["id"])
            for att in atts:
                ext = Path(att["name"]).suffix.lower()
                if ext not in INDEXABLE_EXTENSIONS:
                    print(f"    skip non-indexable: {att['name']}")
                    continue

                # Upload to Drive
                uploaded = upload_file(svc, drive_id, folder_id,
                    att["name"], att["data"], att["contentType"])
                print(f"    uploaded: {att['name']} → {uploaded.get('webViewLink', 'no link')}")

                # Extract + embed
                text = extract_text(att["data"], ext)
                if not text.strip():
                    print(f"    warning: no text extracted from {att['name']}")
                    continue

                title = title_override or Path(att["name"]).stem
                n = embed_and_insert(vo, text, title,
                    f"email-intake/{category}/{att['name']}", category,
                    uploaded["id"], uploaded.get("modifiedTime"),
                    email_message_id=msg_id)
                print(f"    indexed: {n} chunks")
                total_files += 1
                total_chunks += n

        if not email.get("hasAttachments") and len(body_text) <= 50:
            print(f"    no attachments and body too short — skipping")

        # --- Follow URLs found in the email body ---
        urls = extract_urls(body_text)
        if urls:
            print(f"    found {len(urls)} URL(s) to follow")
            for url in urls[:5]:  # cap at 5 URLs per email to avoid runaway
                print(f"    fetching: {url[:100]}")
                if args.dry_run:
                    print(f"      (dry-run — would fetch + convert to PDF)")
                    continue
                pdf_bytes, article_text, page_title = fetch_and_convert_to_pdf(url)
                if not article_text:
                    continue
                safe_title = re.sub(r"[^a-zA-Z0-9\-_ ]", "", page_title or "web_article")[:80]
                pdf_fname = f"{safe_title}.pdf"

                # Upload PDF to Drive if generated successfully
                drive_fid = None
                drive_mod = None
                if pdf_bytes:
                    uploaded_url = upload_file(svc, drive_id, folder_id,
                        pdf_fname, pdf_bytes, "application/pdf")
                    drive_fid = uploaded_url["id"]
                    drive_mod = uploaded_url.get("modifiedTime")
                    print(f"      saved PDF: {pdf_fname} ({len(pdf_bytes)/1024:.0f} KB)")
                else:
                    print(f"      PDF skipped — embedding text only")

                # Embed the extracted text
                url_title = title_override or page_title or safe_title
                n = embed_and_insert(vo, article_text, url_title,
                    f"email-intake/{category}/web/{pdf_fname}", category,
                    drive_fid, drive_mod, email_message_id=msg_id)
                print(f"      indexed: {n} chunks")
                total_files += 1
                total_chunks += n

        # Note: we do NOT mark_read — Mail.Read permission only, no Mail.ReadWrite.
        # Idempotency handled via email_message_id tracking in documents metadata.

    print(f"\n=== Done: {total_files} files, {total_chunks} chunks indexed ===")

    if total_files > 0:
        notify(
            f"CBS KB - {total_files} documents ingested via email",
            f"{total_chunks} chunks embedded from {len(emails)} email(s). Categories: {category}"
        )


if __name__ == "__main__":
    main()
